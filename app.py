from flask import Flask, request, jsonify, render_template
from googletrans import Translator
import os
from werkzeug.utils import secure_filename
import docx
import PyPDF2
from pptx import Presentation

app = Flask(__name__)
translator = Translator()

# Helper function to extract text from PDF
def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text.strip()

# Helper function to extract text from PowerPoint
def extract_text_from_ppt(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# Home route
@app.route("/")
def home():
    return render_template("home.html")

# Document Translation Page
@app.route("/document-translation")
def document_translation():
    return render_template("file-translation.html") 

# API for text translation
@app.route("/translate", methods=["POST"])
def translate_text():
    data = request.json
    text = data.get("text", "").strip()
    target_lang = data.get("target_lang", "hi")

    if not text:
        return jsonify({"error": "Text is empty. Please enter some text."})

    translated = translator.translate(text, dest=target_lang)
    return jsonify({"translated_text": translated.text})

# API for bulk document translation
@app.route("/translate-document", methods=["POST"])
def translate_document():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"})

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"})

    filename = secure_filename(file.filename)
    file_ext = os.path.splitext(filename)[1].lower()

    text = ""
    if file_ext == ".txt":
        text = file.read().decode("utf-8").strip()
    elif file_ext == ".docx":
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs]).strip()
    elif file_ext == ".pdf":
        text = extract_text_from_pdf(file)
    elif file_ext == ".pptx":
        text = extract_text_from_ppt(file)
    else:
        return jsonify({"error": "Unsupported file format"})

    if not text:
        return jsonify({"error": "No text found in the file"})

    target_lang = request.form.get("target_lang", "hi")
    translated = translator.translate(text, dest=target_lang)
    
    return jsonify({"translated_text": translated.text})

@app.route('/languages')
def languages():
    return render_template('languages.html')

if __name__ == "__main__":
    app.run(debug=True)
